"""
Enterprise-Grade Document Processor Factory

Context-Aware RAG chunking with format-specific strategies.
Supports: Code files, Markdown, PDF/DOCX with metadata enrichment.

Author: Axio Hub Team
"""

import io
import os
import re
import logging
import tempfile
import threading
import zipfile
from abc import ABC, abstractmethod
from dataclasses import dataclass, field
from datetime import datetime, timedelta, timezone
from enum import Enum
from typing import List, Dict, Any, Optional, Tuple

# LangChain imports (text_splitter moved to langchain package in newer versions)
# LangChain imports (text_splitter moved to langchain_text_splitters)
from langchain_text_splitters import (
    RecursiveCharacterTextSplitter,
    Language,
    MarkdownHeaderTextSplitter,
)

# Token counting
import tiktoken

logger = logging.getLogger(__name__)
LLAMAPARSE_LOCK = threading.Lock()

# Ghost Protocol: Privacy-safe logging
from core.log_utils import safe_file_ref, safe_extension

try:
    from core.metrics import pdf_scan_detection_total, llamaparse_fallback_total, docx_ocr_fallback_total, image_ocr_total
except Exception:
    pdf_scan_detection_total = None
    llamaparse_fallback_total = None
    docx_ocr_fallback_total = None
    image_ocr_total = None


# =============================================================================
# CIRCUIT BREAKER FOR CLOUD SERVICES
# =============================================================================

class CircuitState(Enum):
    """Circuit breaker states."""
    CLOSED = "closed"       # Normal operation
    OPEN = "open"           # Service blocked
    HALF_OPEN = "half_open" # Testing recovery


class CircuitBreaker:
    """
    Thread-safe circuit breaker for cloud API rate limiting.
    
    - CLOSED: Normal operation, calls go through
    - OPEN: Service blocked, skip directly to fallback
    - HALF_OPEN: After cooldown, try one request to test recovery
    """
    
    def __init__(self, service_name: str, failure_threshold: int = 3, cooldown_seconds: int = 3600):
        self.service_name = service_name
        self.failure_threshold = failure_threshold
        self.cooldown_seconds = cooldown_seconds
        self._state = CircuitState.CLOSED
        self._failure_count = 0
        self._last_failure_time: Optional[datetime] = None
        self._lock = threading.Lock()
    
    def can_execute(self) -> Tuple[bool, str]:
        """Check if request can proceed. Returns (can_execute, reason)."""
        with self._lock:
            if self._state == CircuitState.CLOSED:
                return True, "circuit_closed"
            
            if self._state == CircuitState.OPEN:
                # Check if cooldown has passed
                if self._last_failure_time:
                    elapsed = datetime.now(timezone.utc) - self._last_failure_time
                    if elapsed > timedelta(seconds=self.cooldown_seconds):
                        self._state = CircuitState.HALF_OPEN
                        logger.info(f"[CircuitBreaker] {self.service_name}: HALF_OPEN (testing recovery)")
                        return True, "circuit_half_open_test"
                
                remaining = self.cooldown_seconds
                if self._last_failure_time:
                    remaining = max(0, self.cooldown_seconds - int((datetime.now(timezone.utc) - self._last_failure_time).total_seconds()))
                return False, f"circuit_open_cooldown_{remaining}s"
            
            # HALF_OPEN - allow one test request
            return True, "circuit_half_open_test"
    
    def record_success(self):
        """Record successful call, reset circuit."""
        with self._lock:
            self._failure_count = 0
            self._state = CircuitState.CLOSED
            logger.info(f"[CircuitBreaker] {self.service_name}: CLOSED (success)")
    
    def record_failure(self, error_type: str = "unknown"):
        """Record failed call, potentially open circuit."""
        with self._lock:
            self._failure_count += 1
            self._last_failure_time = datetime.now(timezone.utc)
            
            # Quota errors immediately open the circuit
            if error_type in ("402", "quota_exceeded", "payment_required"):
                self._state = CircuitState.OPEN
                logger.warning(
                    f"[CircuitBreaker] {self.service_name}: OPEN (quota error, "
                    f"blocking for {self.cooldown_seconds}s)"
                )
                return
            
            if self._failure_count >= self.failure_threshold:
                self._state = CircuitState.OPEN
                logger.warning(
                    f"[CircuitBreaker] {self.service_name}: OPEN "
                    f"({self._failure_count} failures, blocking for {self.cooldown_seconds}s)"
                )


# Global circuit breaker for LlamaParse (1 hour cooldown on quota errors)
LLAMAPARSE_CIRCUIT = CircuitBreaker(
    service_name="LlamaParse",
    failure_threshold=3,
    cooldown_seconds=3600
)


# =============================================================================
# CUSTOM EXCEPTIONS
# =============================================================================

class OCRNotAvailableException(Exception):
    """Raised when OCR dependencies are not installed."""
    pass

# Initialize tiktoken encoder (OpenAI's cl100k_base)
try:
    TIKTOKEN_ENCODER = tiktoken.get_encoding("cl100k_base")
except Exception:
    TIKTOKEN_ENCODER = None
    logger.warning("tiktoken encoder not available")


# =============================================================================
# DATA CLASSES
# =============================================================================

@dataclass
class ProcessedChunk:
    """A single processed chunk with content and metadata."""
    content: str
    metadata: Dict[str, Any] = field(default_factory=dict)
    token_count: int = 0
    chunk_index: int = 0


@dataclass
class ProcessedDocument:
    """Result of processing a document."""
    chunks: List[ProcessedChunk]
    file_type: str
    total_tokens: int = 0
    metadata: Dict[str, Any] = field(default_factory=dict)


# =============================================================================
# BASE PROCESSOR
# =============================================================================

class BaseProcessor(ABC):
    """Abstract base class for document processors."""
    
    @abstractmethod
    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        """Process document content and return chunks with metadata."""
        pass
    
    @staticmethod
    def count_tokens(text: str) -> int:
        """Count tokens using tiktoken cl100k_base encoder."""
        if TIKTOKEN_ENCODER:
            return len(TIKTOKEN_ENCODER.encode(text))
        return len(text) // 4  # Fallback approximation


# =============================================================================
# CODE PROCESSOR
# =============================================================================

class CodeProcessor(BaseProcessor):
    """
    Processor for source code files.
    
    Uses RecursiveCharacterTextSplitter.from_language() to preserve
    function and class boundaries. Hard limit of 2000 tokens per chunk.
    """
    
    # Map file extensions to LangChain Language enum
    LANGUAGE_MAP = {
        ".py": Language.PYTHON,
        ".js": Language.JS,
        ".jsx": Language.JS,
        ".ts": Language.TS,
        ".tsx": Language.TS,
        ".java": Language.JAVA,
        ".go": Language.GO,
        ".cpp": Language.CPP,
        ".c": Language.CPP,
        ".cs": Language.CSHARP,
        ".rb": Language.RUBY,
        ".php": Language.PHP,
        ".rs": Language.RUST,
        ".scala": Language.SCALA,
        ".swift": Language.SWIFT,
        ".kt": Language.KOTLIN,
    }
    
    # These don't have special language support, use generic
    GENERIC_CODE_EXTENSIONS = {".json", ".yaml", ".yml", ".toml", ".xml", ".html", ".css", ".sql"}
    
    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        """Process code file with language-aware splitting."""
        ext = os.path.splitext(filename)[1].lower()
        
        # Decode content
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("utf-8", errors="replace")
        
        if not text.strip():
            return ProcessedDocument(chunks=[], file_type="code")
        
        # Get language-specific splitter or generic
        language = self.LANGUAGE_MAP.get(ext)
        
        if language:
            splitter = RecursiveCharacterTextSplitter.from_language(
                language=language,
                chunk_size=1500,  # ~400 tokens target
                chunk_overlap=100,
            )
        else:
            # Generic code splitter for JSON, YAML, etc.
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=1500,
                chunk_overlap=100,
                separators=["\n\n", "\n", " ", ""]
            )
        
        # Split text
        raw_chunks = splitter.split_text(text)
        
        # Build ProcessedChunks with metadata
        chunks = []
        total_tokens = 0
        lang_name = language.value if language else ext.lstrip(".")
        
        for i, chunk_text in enumerate(raw_chunks):
            token_count = self.count_tokens(chunk_text)
            total_tokens += token_count
            
            # Hard limit: if chunk exceeds 2000 tokens, force-split
            if token_count > 2000:
                sub_chunks = self._force_split(chunk_text, 1500)
                for j, sub_text in enumerate(sub_chunks):
                    sub_tokens = self.count_tokens(sub_text)
                    chunks.append(ProcessedChunk(
                        content=sub_text,
                        metadata={
                            "file_type": "code",
                            "language": lang_name,
                            "filename": filename,
                        },
                        token_count=sub_tokens,
                        chunk_index=len(chunks)
                    ))
            else:
                chunks.append(ProcessedChunk(
                    content=chunk_text,
                    metadata={
                        "file_type": "code",
                        "language": lang_name,
                        "filename": filename,
                    },
                    token_count=token_count,
                    chunk_index=i
                ))
        
        # Re-index after potential sub-splits
        for i, chunk in enumerate(chunks):
            chunk.chunk_index = i
        
        logger.info(f"[CodeProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens")
        return ProcessedDocument(
            chunks=chunks,
            file_type="code",
            total_tokens=total_tokens,
            metadata={"language": lang_name}
        )
    
    def _force_split(self, text: str, max_chars: int) -> List[str]:
        """Force-split oversized text by character count."""
        return [text[i:i+max_chars] for i in range(0, len(text), max_chars)]


# =============================================================================
# MARKDOWN PROCESSOR
# =============================================================================

class MarkdownProcessor(BaseProcessor):
    """
    Processor for Markdown files and web content.
    
    Strategy:
    1. Split by headers (#, ##, ###) first
    2. For each section, apply recursive character splitting
    3. Inject header path as context prefix
    """
    
    HEADERS_TO_SPLIT_ON = [
        ("#", "Header1"),
        ("##", "Header2"),
        ("###", "Header3"),
        ("####", "Header4"),
    ]
    
    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        """Process markdown with header-aware splitting."""
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("utf-8", errors="replace")
        
        if not text.strip():
            return ProcessedDocument(chunks=[], file_type="markdown")
        
        # Step 1: Split by headers
        header_splitter = MarkdownHeaderTextSplitter(
            headers_to_split_on=self.HEADERS_TO_SPLIT_ON,
            strip_headers=False
        )
        
        try:
            header_docs = header_splitter.split_text(text)
        except Exception as e:
            logger.warning(f"[MarkdownProcessor] Header splitting failed: {e}, using fallback")
            header_docs = None
        
        chunks = []
        total_tokens = 0
        
        if header_docs:
            # Process each header section
            content_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=150,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
            
            for doc in header_docs:
                # Build header path from metadata
                header_path = self._build_header_path(doc.metadata)
                section_text = doc.page_content
                
                # Split section content
                section_chunks = content_splitter.split_text(section_text)
                
                for chunk_text in section_chunks:
                    # Context injection: prepend header path
                    if header_path:
                        contextualized = f"[Context: {header_path}]\n{chunk_text}"
                    else:
                        contextualized = chunk_text
                    
                    token_count = self.count_tokens(contextualized)
                    total_tokens += token_count
                    
                    chunks.append(ProcessedChunk(
                        content=contextualized,
                        metadata={
                            "file_type": "markdown",
                            "header_path": header_path,
                            "filename": filename,
                        },
                        token_count=token_count,
                        chunk_index=len(chunks)
                    ))
        else:
            # Fallback: simple recursive split
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=150
            )
            raw_chunks = splitter.split_text(text)
            
            for i, chunk_text in enumerate(raw_chunks):
                token_count = self.count_tokens(chunk_text)
                total_tokens += token_count
                
                chunks.append(ProcessedChunk(
                    content=chunk_text,
                    metadata={
                        "file_type": "markdown",
                        "header_path": "",
                        "filename": filename,
                    },
                    token_count=token_count,
                    chunk_index=i
                ))
        
        logger.info(f"[MarkdownProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens")
        return ProcessedDocument(
            chunks=chunks,
            file_type="markdown",
            total_tokens=total_tokens
        )
    
    def _build_header_path(self, metadata: Dict[str, Any]) -> str:
        """Build header path string from metadata."""
        parts = []
        for key in ["Header1", "Header2", "Header3", "Header4"]:
            if key in metadata and metadata[key]:
                # Clean header text (remove # symbols)
                header = metadata[key].strip().lstrip("#").strip()
                if header:
                    parts.append(header)
        return " > ".join(parts)


# =============================================================================
# PDF PROCESSOR
# =============================================================================

class PDFProcessor(BaseProcessor):
    """
    Enterprise PDF Processor with 3-Tier Resilient Architecture.
    
    Tier 1 (Premium Cloud): LlamaParse - OCR, tables, complex layouts
    Tier 2 (Fast Local): PyMuPDF - Text layer extraction
    Tier 3 (Smart Local): Tesseract OCR - Scanned document handling
    
    Features:
    - Circuit breaker for cloud quota management
    - Automatic scanned PDF detection
    - Zero cloud dependency fallback
    """
    
    # Regex patterns for common headers/footers to remove
    NOISE_PATTERNS = [
        r"Page\s+\d+\s+(of|/)\s+\d+",  # Page 1 of 10
        r"^\d+\s*$",  # Lone page numbers
        r"CONFIDENTIAL",
        r"^\s*©.*$",  # Copyright notices
    ]
    
    # Thresholds for content validation
    SCANNED_TEXT_THRESHOLD = 150      # chars - below this likely scanned
    MIN_TOKENS_THRESHOLD = 50         # tokens - below this trigger OCR fallback
    OCR_QUALITY_THRESHOLD = 100       # tokens - minimum acceptable OCR result
    
    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        """
        Process PDF with 3-tier cascading fallback.
        
        Flow:
        1. Check circuit breaker for LlamaParse
        2. Try LlamaParse if available (handles 402 gracefully)
        3. Try PyMuPDF for text extraction
        4. If low content, try Tesseract OCR
        """
        from core.config import settings
        
        # =====================================================================
        # TIER 1: LlamaParse (Cloud Premium)
        # =====================================================================
        if settings.LLAMA_CLOUD_API_KEY:
            can_execute, reason = LLAMAPARSE_CIRCUIT.can_execute()
            
            if can_execute:
                logger.info(f"[PDFProcessor] Tier 1: Trying LlamaParse for {safe_file_ref(filename=filename)}")
                try:
                    result = self._process_with_llamaparse(content, filename)
                    if result and result.chunks and result.total_tokens >= self.MIN_TOKENS_THRESHOLD:
                        LLAMAPARSE_CIRCUIT.record_success()
                        if pdf_scan_detection_total:
                            pdf_scan_detection_total.labels("tier1_llamaparse_success").inc()
                        return result
                    
                    logger.warning(
                        f"[PDFProcessor] LlamaParse returned insufficient content "
                        f"({result.total_tokens if result else 0} tokens)"
                    )
                    
                except Exception as e:
                    error_str = str(e).lower()
                    # Check for quota/payment errors
                    if "402" in error_str or "payment" in error_str or "quota" in error_str or "credits" in error_str:
                        LLAMAPARSE_CIRCUIT.record_failure("402")
                        logger.warning(f"[PDFProcessor] LlamaParse quota exceeded: {e}")
                    else:
                        LLAMAPARSE_CIRCUIT.record_failure("error")
                        logger.warning(f"[PDFProcessor] LlamaParse failed: {e}")
            else:
                logger.info(f"[PDFProcessor] Skipping LlamaParse ({reason})")
                if llamaparse_fallback_total:
                    llamaparse_fallback_total.labels(reason).inc()
        else:
            logger.info(f"[PDFProcessor] No LLAMA_CLOUD_API_KEY, using local processing for {safe_file_ref(filename=filename)}")
            if pdf_scan_detection_total:
                pdf_scan_detection_total.labels("local_no_api_key").inc()
        
        # =====================================================================
        # TIER 2: PyMuPDF (Fast Local)
        # =====================================================================
        logger.info(f"[PDFProcessor] Tier 2: Trying PyMuPDF for {safe_file_ref(filename=filename)}")
        result = self._process_with_pymupdf(content, filename)
        
        if result and result.chunks and result.total_tokens >= self.MIN_TOKENS_THRESHOLD:
            if pdf_scan_detection_total:
                pdf_scan_detection_total.labels("tier2_pymupdf_success").inc()
            return result
        
        logger.warning(
            f"[PDFProcessor] PyMuPDF low content ({result.total_tokens if result else 0} tokens), "
            f"likely scanned PDF - triggering OCR"
        )
        
        # =====================================================================
        # TIER 3: Tesseract OCR (Smart Local Fallback)
        # =====================================================================
        logger.info(f"[PDFProcessor] Tier 3: Trying Tesseract OCR for {safe_file_ref(filename=filename)}")
        try:
            ocr_result = self._process_with_ocr(content, filename)
            if ocr_result and ocr_result.chunks and ocr_result.total_tokens >= self.OCR_QUALITY_THRESHOLD:
                if pdf_scan_detection_total:
                    pdf_scan_detection_total.labels("tier3_ocr_success").inc()
                return ocr_result
            
            logger.warning(f"[PDFProcessor] OCR also returned low content for {safe_file_ref(filename=filename)}")
            
        except OCRNotAvailableException as e:
            logger.warning(f"[PDFProcessor] OCR not available: {e}")
        except Exception as e:
            logger.error(f"[PDFProcessor] OCR failed: {e}")
        
        # Return best available result (even if low quality)
        if pdf_scan_detection_total:
            pdf_scan_detection_total.labels("all_tiers_low_content").inc()
        
        return result or ProcessedDocument(chunks=[], file_type="pdf")

    def _is_likely_scanned(self, text_length: int) -> bool:
        """Heuristic: very low extracted text likely indicates scanned PDF."""
        return text_length < self.SCANNED_TEXT_THRESHOLD
    
    def _process_with_llamaparse(self, content: bytes, filename: str) -> ProcessedDocument:
        """
        Process PDF using LlamaParse for advanced OCR and table extraction.
        
        LlamaParse provides:
        - OCR for scanned documents
        - Table detection and extraction
        - Better handling of complex layouts
        """
        import tempfile
        import os
        
        # Import LlamaParse with nest_asyncio for async compatibility
        try:
            import nest_asyncio
            nest_asyncio.apply()
            from llama_parse import LlamaParse
        except ImportError:
            raise ImportError("llama-parse package not installed")
        
        from core.config import settings
        
        logger.info(f"[PDFProcessor] Using LlamaParse for {safe_file_ref(filename=filename)}")
        
        # Write content to temp file (LlamaParse needs file path)
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tf:
            tf.write(content)
            tf_path = tf.name
        
        try:
            # Initialize LlamaParse with OCR settings
            # NOTE: Serialize LlamaParse calls to avoid anyio/gevent conflicts.
            with LLAMAPARSE_LOCK:
                parser = LlamaParse(
                    api_key=settings.LLAMA_CLOUD_API_KEY,
                    result_type="markdown",  # Get structured markdown output
                    verbose=False,
                    language="en",
                )
                
                # Parse the document (synchronous call)
                documents = parser.load_data(tf_path)
            
            if not documents:
                raise ValueError("LlamaParse returned no documents")
            
            # Combine all document text
            full_text = "\n\n".join([doc.text for doc in documents])
            
            if not full_text.strip():
                raise ValueError("LlamaParse returned empty text")
            
            logger.info(f"[PDFProcessor] LlamaParse extracted {len(full_text)} chars from {safe_file_ref(filename=filename)}")
            
        finally:
            # Always cleanup temp file
            try:
                os.remove(tf_path)
            except:
                pass
        
        # Chunk the extracted text (use MarkdownProcessor for markdown output)
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        raw_chunks = splitter.split_text(full_text)
        
        chunks = []
        total_tokens = 0
        for i, chunk_text in enumerate(raw_chunks):
            contextualized = f"[File: {filename}]\n{chunk_text}"
            token_count = self.count_tokens(contextualized)
            total_tokens += token_count
            
            chunks.append(ProcessedChunk(
                content=contextualized,
                metadata={
                    "file_type": "pdf",
                    "parser": "llama_parse",
                    "filename": filename,
                },
                token_count=token_count,
                chunk_index=i
            ))
        
        logger.info(f"[PDFProcessor] LlamaParse: {safe_file_ref(filename=filename)}: {len(chunks)} chunks")
        return ProcessedDocument(
            chunks=chunks,
            file_type="pdf",
            total_tokens=total_tokens,
            metadata={"parser": "llama_parse"}
        )
    
    def _process_with_pymupdf(self, content: bytes, filename: str) -> ProcessedDocument:
        """Process PDF with PyMuPDF (standard mode)."""
        try:
            import fitz  # PyMuPDF
        except ImportError:
            logger.error("[PDFProcessor] PyMuPDF not installed, using fallback")
            return self._fallback_process(content, filename)
        
        # Extract text from PDF
        pages_text = []
        total_chars = 0
        try:
            doc = fitz.open(stream=content, filetype="pdf")
            for page_num, page in enumerate(doc, start=1):
                text = page.get_text("text")
                if text.strip():
                    # Clean the text
                    cleaned = self._clean_text(text)
                    if cleaned.strip():
                        pages_text.append((page_num, cleaned))
                        total_chars += len(cleaned)
            doc.close()
        except Exception as e:
            logger.error(f"[PDFProcessor] PyMuPDF extraction failed: {e}")
            return self._fallback_process(content, filename)
        
        if not pages_text:
            return ProcessedDocument(chunks=[], file_type="pdf")
        
        # Chunk with sliding window
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        
        chunks = []
        total_tokens = 0
        
        for page_num, page_text in pages_text:
            page_chunks = splitter.split_text(page_text)
            
            for chunk_text in page_chunks:
                # Context injection: prepend file and page info
                contextualized = f"[File: {filename}] [Page: {page_num}]\n{chunk_text}"
                
                token_count = self.count_tokens(contextualized)
                total_tokens += token_count
                
                chunks.append(ProcessedChunk(
                    content=contextualized,
                    metadata={
                        "file_type": "pdf",
                        "parser": "pymupdf",
                        "page_number": page_num,
                        "filename": filename,
                    },
                    token_count=token_count,
                    chunk_index=len(chunks)
                ))
        
        logger.info(f"[PDFProcessor] PyMuPDF: {safe_file_ref(filename=filename)}: {len(chunks)} chunks from {len(pages_text)} pages")
        return ProcessedDocument(
            chunks=chunks,
            file_type="pdf",
            total_tokens=total_tokens,
            metadata={"total_pages": len(pages_text), "parser": "pymupdf", "text_length": total_chars}
        )
    
    def _clean_text(self, text: str) -> str:
        """Remove common headers, footers, and noise patterns."""
        lines = text.split("\n")
        cleaned_lines = []
        
        for line in lines:
            skip = False
            for pattern in self.NOISE_PATTERNS:
                if re.search(pattern, line, re.IGNORECASE):
                    skip = True
                    break
            if not skip:
                cleaned_lines.append(line)
        
        return "\n".join(cleaned_lines)
    
    def _fallback_process(self, content: bytes, filename: str) -> ProcessedDocument:
        """Fallback using pypdf if PyMuPDF is not available."""
        try:
            from pypdf import PdfReader
            reader = PdfReader(io.BytesIO(content))
            full_text = ""
            for page in reader.pages:
                text = page.extract_text() or ""
                full_text += text + "\n\n"
        except Exception as e:
            logger.error(f"[PDFProcessor] pypdf fallback failed: {e}")
            return ProcessedDocument(chunks=[], file_type="pdf")
        
        # Simple chunking
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200
        )
        raw_chunks = splitter.split_text(full_text)
        
        chunks = []
        total_tokens = 0
        for i, chunk_text in enumerate(raw_chunks):
            contextualized = f"[File: {filename}]\n{chunk_text}"
            token_count = self.count_tokens(contextualized)
            total_tokens += token_count
            chunks.append(ProcessedChunk(
                content=contextualized,
                metadata={"file_type": "pdf", "filename": filename},
                token_count=token_count,
                chunk_index=i
            ))
        
        return ProcessedDocument(chunks=chunks, file_type="pdf", total_tokens=total_tokens)
    
    def _process_with_ocr(self, content: bytes, filename: str) -> ProcessedDocument:
        """
        Tier 3: Tesseract OCR for scanned PDFs.
        
        Converts PDF pages to images and runs OCR on each page.
        Requires: tesseract-ocr, poppler-utils (system), pytesseract, pdf2image (Python)
        """
        try:
            from pdf2image import convert_from_bytes
            import pytesseract
        except ImportError as e:
            raise OCRNotAvailableException(f"OCR dependencies not installed: {e}")
        
        logger.info(f"[PDFProcessor] Running Tesseract OCR on {safe_file_ref(filename=filename)}")
        
        try:
            # Convert PDF pages to images (300 DPI for good OCR quality)
            images = convert_from_bytes(content, dpi=300, fmt='png')
            
            ocr_texts = []
            for page_num, image in enumerate(images, start=1):
                # Run OCR on each page
                text = pytesseract.image_to_string(image, lang='eng')
                if text.strip():
                    cleaned = self._clean_text(text)
                    if cleaned.strip():
                        ocr_texts.append(f"[Page {page_num}]\n{cleaned}")
            
            if not ocr_texts:
                return ProcessedDocument(chunks=[], file_type="pdf", total_tokens=0)
            
            full_text = "\n\n".join(ocr_texts)
            logger.info(f"[PDFProcessor] OCR: {safe_file_ref(filename=filename)}: {len(full_text)} chars from {len(images)} pages")
            
            # Chunk the OCR text
            splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200,
                separators=["\n\n", "\n", ". ", " ", ""]
            )
            raw_chunks = splitter.split_text(full_text)
            
            chunks = []
            total_tokens = 0
            for i, chunk_text in enumerate(raw_chunks):
                contextualized = f"[File: {filename}]\n{chunk_text}"
                token_count = self.count_tokens(contextualized)
                total_tokens += token_count
                
                chunks.append(ProcessedChunk(
                    content=contextualized,
                    metadata={
                        "file_type": "pdf",
                        "parser": "tesseract_ocr",
                        "filename": filename,
                    },
                    token_count=token_count,
                    chunk_index=i
                ))
            
            logger.info(f"[PDFProcessor] Tesseract OCR: {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens")
            return ProcessedDocument(
                chunks=chunks,
                file_type="pdf",
                total_tokens=total_tokens,
                metadata={"parser": "tesseract_ocr", "pages_ocr": len(images)}
            )
            
        except Exception as e:
            logger.error(f"[PDFProcessor] OCR processing failed: {e}")
            raise


# =============================================================================
# DOCX PROCESSOR (ENHANCED WITH OCR FALLBACK)
# =============================================================================

class DocxProcessor(BaseProcessor):
    """
    Processor for Word documents with 3-Tier Resilient Architecture.
    
    Tier 1 (Fast Local): docx2txt - Text layer extraction
    Tier 2 (OCR Fallback): Extract embedded images from DOCX and run Tesseract
    Tier 3 (Cloud Premium): LlamaParse - Advanced OCR with table detection
    
    DOCX files are actually ZIP archives containing:
    - word/document.xml - main content
    - word/media/ - embedded images (image1.png, image2.jpeg, etc.)
    
    This handles scanned contracts, signed documents, and image-heavy DOCXs.
    """
    
    MIN_TOKENS_THRESHOLD = 50   # Below this, trigger OCR fallback
    OCR_QUALITY_THRESHOLD = 30  # Minimum acceptable OCR result
    
    # Supported image extensions in DOCX media folder
    IMAGE_EXTENSIONS = {'.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.tif', '.emf', '.wmf'}
    
    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        """
        Process DOCX with 3-tier cascading fallback.
        
        Flow:
        1. Try docx2txt for text extraction
        2. If low content, extract embedded images and OCR them
        3. If still low, try LlamaParse (if configured)
        """
        
        # =====================================================================
        # TIER 1: docx2txt (Fast Local)
        # =====================================================================
        logger.info(f"[DocxProcessor] Tier 1: Trying docx2txt for {safe_file_ref(filename=filename)}")
        text = ""
        try:
            import docx2txt
            text = docx2txt.process(io.BytesIO(content)) or ""
        except Exception as e:
            logger.warning(f"[DocxProcessor] docx2txt extraction failed: {e}")
        
        # Check if we have sufficient content
        if text.strip():
            token_count = self.count_tokens(text)
            if token_count >= self.MIN_TOKENS_THRESHOLD:
                logger.info(f"[DocxProcessor] Tier 1 success: {token_count} tokens extracted from {safe_file_ref(filename=filename)}")
                if docx_ocr_fallback_total:
                    docx_ocr_fallback_total.labels("tier1_docx2txt_success").inc()
                return self._chunk_text(text, filename, parser="docx2txt")
            else:
                logger.warning(
                    f"[DocxProcessor] Low content ({token_count} tokens) in {filename}, "
                    f"likely scanned DOCX - triggering OCR"
                )
        else:
            logger.warning(f"[DocxProcessor] No text extracted from {safe_file_ref(filename=filename)}, likely scanned DOCX - triggering OCR")
        
        # =====================================================================
        # TIER 2: Extract embedded images and OCR them
        # =====================================================================
        logger.info(f"[DocxProcessor] Tier 2: Trying image extraction + OCR for {safe_file_ref(filename=filename)}")
        try:
            ocr_result = self._process_with_embedded_image_ocr(content, filename)
            if ocr_result and ocr_result.chunks and ocr_result.total_tokens >= self.OCR_QUALITY_THRESHOLD:
                logger.info(f"[DocxProcessor] Tier 2 OCR success: {ocr_result.total_tokens} tokens from {safe_file_ref(filename=filename)}")
                if docx_ocr_fallback_total:
                    docx_ocr_fallback_total.labels("tier2_tesseract_success").inc()
                return ocr_result
            logger.warning(f"[DocxProcessor] OCR returned low content ({ocr_result.total_tokens if ocr_result else 0} tokens) for {safe_file_ref(filename=filename)}")
        except OCRNotAvailableException as e:
            logger.warning(f"[DocxProcessor] OCR not available: {e}")
        except Exception as e:
            logger.error(f"[DocxProcessor] OCR failed for {safe_file_ref(filename=filename)}: {e}")
        
        # =====================================================================
        # TIER 3: LlamaParse (Cloud Premium Fallback)
        # =====================================================================
        from core.config import settings
        if settings.LLAMA_CLOUD_API_KEY:
            can_execute, reason = LLAMAPARSE_CIRCUIT.can_execute()
            if can_execute:
                logger.info(f"[DocxProcessor] Tier 3: Trying LlamaParse for {safe_file_ref(filename=filename)}")
                try:
                    result = LlamaParseProcessor(file_type="docx").process(content, filename)
                    if result and result.chunks and result.total_tokens >= self.OCR_QUALITY_THRESHOLD:
                        LLAMAPARSE_CIRCUIT.record_success()
                        if docx_ocr_fallback_total:
                            docx_ocr_fallback_total.labels("tier3_llamaparse_success").inc()
                        logger.info(f"[DocxProcessor] Tier 3 LlamaParse success: {result.total_tokens} tokens from {safe_file_ref(filename=filename)}")
                        return result
                except Exception as e:
                    error_str = str(e).lower()
                    if "402" in error_str or "quota" in error_str or "payment" in error_str or "credits" in error_str:
                        LLAMAPARSE_CIRCUIT.record_failure("402")
                        logger.warning(f"[DocxProcessor] LlamaParse quota exceeded: {e}")
                    else:
                        LLAMAPARSE_CIRCUIT.record_failure("error")
                        logger.warning(f"[DocxProcessor] LlamaParse failed: {e}")
            else:
                logger.info(f"[DocxProcessor] Skipping LlamaParse ({reason}) for {safe_file_ref(filename=filename)}")
                if llamaparse_fallback_total:
                    llamaparse_fallback_total.labels(f"docx_{reason}").inc()
        
        # Return best available result (even if low quality)
        if docx_ocr_fallback_total:
            docx_ocr_fallback_total.labels("all_tiers_low_content").inc()
        
        if text.strip():
            logger.warning(f"[DocxProcessor] Returning low-content result for {safe_file_ref(filename=filename)}")
            return self._chunk_text(text, filename, parser="docx2txt_low_content")
        
        logger.error(f"[DocxProcessor] All tiers failed for {safe_file_ref(filename=filename)}, returning empty result")
        return ProcessedDocument(chunks=[], file_type="docx")
        
    def _process_with_embedded_image_ocr(self, content: bytes, filename: str) -> ProcessedDocument:
        """
        Extract embedded images from DOCX (which is a ZIP file) and OCR them.
        
        DOCX structure:
        - word/media/image1.png
        - word/media/image2.jpeg
        - etc.
        
        Returns:
            ProcessedDocument with OCR'd text from embedded images
        """
        try:
            import pytesseract
            from PIL import Image
        except ImportError as e:
            raise OCRNotAvailableException(f"OCR dependencies not installed: {e}")
        
        ocr_texts = []
        images_processed = 0
        images_failed = 0
        
        try:
            with zipfile.ZipFile(io.BytesIO(content), 'r') as zf:
                # Find all image files in word/media/
                image_files = [
                    name for name in zf.namelist()
                    if name.startswith('word/media/') and 
                    os.path.splitext(name)[1].lower() in self.IMAGE_EXTENSIONS
                ]
                
                if not image_files:
                    logger.info(f"[DocxProcessor] No embedded images found in {safe_file_ref(filename=filename)}")
                    return ProcessedDocument(chunks=[], file_type="docx", total_tokens=0)
                
                logger.info(f"[DocxProcessor] Found {len(image_files)} embedded images in {safe_file_ref(filename=filename)}")
                
                # Sort to maintain order (image1, image2, etc.)
                image_files.sort()
                
                for idx, img_path in enumerate(image_files, start=1):
                    try:
                        # Extract image from ZIP
                        img_data = zf.read(img_path)
                        image = Image.open(io.BytesIO(img_data))
                        
                        # Skip very small images (likely icons/bullets)
                        if image.width < 50 or image.height < 50:
                            logger.debug(f"[DocxProcessor] Skipping small image {img_path} ({image.width}x{image.height})")
                            continue
                        
                        # Convert to RGB if necessary (for RGBA/palette/grayscale images)
                        if image.mode in ('RGBA', 'P', 'LA', 'L'):
                            # Create white background for transparent images
                            if image.mode in ('RGBA', 'LA', 'P'):
                                background = Image.new('RGB', image.size, (255, 255, 255))
                                if image.mode == 'P':
                                    image = image.convert('RGBA')
                                background.paste(image, mask=image.split()[-1] if image.mode == 'RGBA' else None)
                                image = background
                            else:
                                image = image.convert('RGB')
                        
                        # Run OCR with language detection hint
                        # Use multiple languages for better accuracy
                        try:
                            text = pytesseract.image_to_string(image, lang='eng+tur')
                        except Exception:
                            # Fallback to English only if multi-lang fails
                            text = pytesseract.image_to_string(image, lang='eng')
                        
                        if text.strip():
                            # Clean up OCR artifacts
                            cleaned_text = self._clean_ocr_text(text)
                            if cleaned_text.strip():
                                ocr_texts.append(f"[Embedded Image {idx}]\n{cleaned_text.strip()}")
                                images_processed += 1
                                
                    except Exception as e:
                        images_failed += 1
                        logger.warning(f"[DocxProcessor] Failed to OCR image {img_path}: {e}")
                        continue
                        
        except zipfile.BadZipFile:
            logger.error(f"[DocxProcessor] Invalid DOCX (not a valid ZIP): {safe_file_ref(filename=filename)}")
            raise OCRNotAvailableException("Invalid DOCX format - not a valid ZIP archive")
        except Exception as e:
            logger.error(f"[DocxProcessor] Failed to process DOCX as ZIP: {e}")
            raise OCRNotAvailableException(f"Failed to extract images from DOCX: {e}")
        
        if not ocr_texts:
            logger.info(f"[DocxProcessor] No text extracted from {len(image_files)} images in {safe_file_ref(filename=filename)}")
            return ProcessedDocument(chunks=[], file_type="docx", total_tokens=0)
        
        full_text = "\n\n".join(ocr_texts)
        logger.info(
            f"[DocxProcessor] OCR extracted {len(full_text)} chars from "
            f"{images_processed}/{len(image_files)} images in {filename} "
            f"({images_failed} failed)"
        )
        
        return self._chunk_text(full_text, filename, parser="tesseract_ocr")
    
    def _clean_ocr_text(self, text: str) -> str:
        """Clean up common OCR artifacts and noise."""
        # Remove excessive whitespace
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r' {2,}', ' ', text)
        
        # Remove common OCR noise patterns
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            # Skip lines that are just noise (single characters, numbers only, etc.)
            if len(line) <= 2 and not line.isalnum():
                continue
            # Skip lines that are mostly special characters
            if line and sum(1 for c in line if c.isalnum()) / len(line) < 0.3:
                continue
            cleaned_lines.append(line)
        
        return '\n'.join(cleaned_lines)
    
    def _chunk_text(self, text: str, filename: str, parser: str = "docx2txt") -> ProcessedDocument:
        """Chunk extracted text into ProcessedChunks with context injection."""
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        raw_chunks = splitter.split_text(text)
        
        chunks = []
        total_tokens = 0
        for i, chunk_text in enumerate(raw_chunks):
            contextualized = f"[File: {filename}]\n{chunk_text}"
            token_count = self.count_tokens(contextualized)
            total_tokens += token_count
            chunks.append(ProcessedChunk(
                content=contextualized,
                metadata={
                    "file_type": "docx",
                    "parser": parser,
                    "filename": filename
                },
                token_count=token_count,
                chunk_index=i
            ))
        
        logger.info(f"[DocxProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens (parser={parser})")
        return ProcessedDocument(
            chunks=chunks,
            file_type="docx",
            total_tokens=total_tokens,
            metadata={"parser": parser}
        )


# =============================================================================
# HTML PROCESSOR
# =============================================================================

class HTMLProcessor(BaseProcessor):
    """Processor for HTML files (strip tags, keep readable text)."""

    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("utf-8", errors="replace")

        if not text.strip():
            return ProcessedDocument(chunks=[], file_type="html")

        try:
            from bs4 import BeautifulSoup
        except ImportError as exc:
            logger.error(f"[HTMLProcessor] Missing dependency: {exc}")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "missing_dependency"},
            )

        soup = BeautifulSoup(text, "html.parser")
        extracted = soup.get_text(separator="\n")
        extracted = re.sub(r"\n{3,}", "\n\n", extracted).strip()

        if not extracted:
            return ProcessedDocument(chunks=[], file_type="html")

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=150,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        raw_chunks = splitter.split_text(extracted)

        chunks = []
        total_tokens = 0
        for i, chunk_text in enumerate(raw_chunks):
            token_count = self.count_tokens(chunk_text)
            total_tokens += token_count
            chunks.append(
                ProcessedChunk(
                    content=chunk_text,
                    metadata={"file_type": "html", "filename": filename},
                    token_count=token_count,
                    chunk_index=i,
                )
            )

        logger.info(f"[HTMLProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks")
        return ProcessedDocument(chunks=chunks, file_type="html", total_tokens=total_tokens)


# =============================================================================
# CSV PROCESSOR
# =============================================================================

class CSVProcessor(BaseProcessor):
    """Processor for CSV/TSV files with structured row output."""

    CHUNK_SIZE = 500

    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        from core.config import settings

        if len(content) > settings.MAX_STRUCTURED_FILE_SIZE:
            logger.warning(f"[CSVProcessor] File too large for CSV parsing: {safe_file_ref(filename=filename)}")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "file_too_large"},
            )

        try:
            import pandas as pd
        except ImportError as exc:
            logger.error(f"[CSVProcessor] Missing dependency: {exc}")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "missing_dependency"},
            )

        delimiter = "\t" if filename.lower().endswith(".tsv") else ","
        chunks = []
        total_tokens = 0

        try:
            df_iter = pd.read_csv(
                io.BytesIO(content),
                chunksize=self.CHUNK_SIZE,
                encoding="utf-8",
                on_bad_lines="skip",
                sep=delimiter,
            )

            for chunk_idx, df_chunk in enumerate(df_iter):
                text_rows = []
                for _idx, row in df_chunk.iterrows():
                    row_text = " | ".join(
                        f"{col}: {val}"
                        for col, val in row.items()
                        if pd.notna(val)
                    )
                    if row_text:
                        text_rows.append(row_text)

                if not text_rows:
                    continue

                start_row = chunk_idx * self.CHUNK_SIZE + 1
                end_row = start_row + len(df_chunk) - 1
                chunk_text = (
                    f"[File: {filename}] [Rows: {start_row}-{end_row}]\n"
                    + "\n".join(text_rows)
                )

                token_count = self.count_tokens(chunk_text)
                total_tokens += token_count
                chunks.append(
                    ProcessedChunk(
                        content=chunk_text,
                        metadata={
                            "file_type": "csv",
                            "row_range": f"{start_row}-{end_row}",
                            "filename": filename,
                        },
                        token_count=token_count,
                        chunk_index=len(chunks),
                    )
                )

        except Exception as exc:
            logger.error(f"[CSVProcessor] Failed to parse {safe_file_ref(filename=filename)}: {exc}")
            return ProcessedDocument(chunks=[], file_type="csv")

        logger.info(f"[CSVProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens")
        return ProcessedDocument(chunks=chunks, file_type="csv", total_tokens=total_tokens)


# =============================================================================
# EXCEL PROCESSOR
# =============================================================================

class ExcelProcessor(BaseProcessor):
    """Processor for XLSX files using streaming read_only mode."""

    CHUNK_SIZE = 200

    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        from core.config import settings

        if len(content) > settings.MAX_STRUCTURED_FILE_SIZE:
            logger.warning(f"[ExcelProcessor] File too large for XLSX parsing: {safe_file_ref(filename=filename)}")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "file_too_large"},
            )

        try:
            import openpyxl
        except ImportError as exc:
            logger.error(f"[ExcelProcessor] Missing dependency: {exc}")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "missing_dependency"},
            )

        chunks = []
        total_tokens = 0

        try:
            workbook = openpyxl.load_workbook(
                io.BytesIO(content),
                read_only=True,
                data_only=True,
            )
        except Exception as exc:
            logger.error(f"[ExcelProcessor] Failed to load workbook {safe_file_ref(filename=filename)}: {exc}")
            return ProcessedDocument(chunks=[], file_type="xlsx")

        try:
            for sheet in workbook.worksheets:
                row_iter = sheet.iter_rows(values_only=True)
                header = None
                buffer = []
                start_row = 1
                current_row = 0

                for row in row_iter:
                    current_row += 1
                    if header is None:
                        header = [
                            (str(cell).strip() if cell is not None else f"Column {idx + 1}")
                            for idx, cell in enumerate(row)
                        ]
                        start_row = current_row + 1
                        continue

                    values = []
                    for idx, cell in enumerate(row):
                        if cell is None:
                            continue
                        label = header[idx] if idx < len(header) else f"Column {idx + 1}"
                        values.append(f"{label}: {cell}")

                    if values:
                        buffer.append(" | ".join(values))

                    if len(buffer) >= self.CHUNK_SIZE:
                        end_row = current_row
                        chunk_text = (
                            f"[File: {filename}] [Sheet: {sheet.title}] [Rows: {start_row}-{end_row}]\n"
                            + "\n".join(buffer)
                        )
                        token_count = self.count_tokens(chunk_text)
                        total_tokens += token_count
                        chunks.append(
                            ProcessedChunk(
                                content=chunk_text,
                                metadata={
                                    "file_type": "xlsx",
                                    "sheet": sheet.title,
                                    "row_range": f"{start_row}-{end_row}",
                                    "filename": filename,
                                },
                                token_count=token_count,
                                chunk_index=len(chunks),
                            )
                        )
                        buffer = []
                        start_row = end_row + 1

                if buffer:
                    end_row = current_row
                    chunk_text = (
                        f"[File: {filename}] [Sheet: {sheet.title}] [Rows: {start_row}-{end_row}]\n"
                        + "\n".join(buffer)
                    )
                    token_count = self.count_tokens(chunk_text)
                    total_tokens += token_count
                    chunks.append(
                        ProcessedChunk(
                            content=chunk_text,
                            metadata={
                                "file_type": "xlsx",
                                "sheet": sheet.title,
                                "row_range": f"{start_row}-{end_row}",
                                "filename": filename,
                            },
                            token_count=token_count,
                            chunk_index=len(chunks),
                        )
                    )

        finally:
            try:
                workbook.close()
            except Exception:
                pass

        logger.info(f"[ExcelProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens")
        return ProcessedDocument(chunks=chunks, file_type="xlsx", total_tokens=total_tokens)


# =============================================================================
# PPTX PROCESSOR
# =============================================================================

class PPTXProcessor(BaseProcessor):
    """Processor for PPTX files using python-pptx."""

    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        try:
            from pptx import Presentation
        except ImportError as exc:
            logger.error(f"[PPTXProcessor] Missing dependency: {exc}")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "missing_dependency"},
            )

        chunks = []
        total_tokens = 0

        try:
            presentation = Presentation(io.BytesIO(content))
        except Exception as exc:
            logger.error(f"[PPTXProcessor] Failed to load {safe_file_ref(filename=filename)}: {exc}")
            return ProcessedDocument(chunks=[], file_type="pptx")

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=150,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_texts = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text or ""
                    if text.strip():
                        slide_texts.append(text.strip())

            if getattr(slide, "has_notes_slide", False) and slide.notes_slide:
                notes = slide.notes_slide.notes_text_frame.text or ""
                if notes.strip():
                    slide_texts.append(f"[Notes]\n{notes.strip()}")

            if not slide_texts:
                continue

            slide_text = "\n".join(slide_texts)
            for chunk_text in splitter.split_text(slide_text):
                token_count = self.count_tokens(chunk_text)
                total_tokens += token_count
                chunks.append(
                    ProcessedChunk(
                        content=chunk_text,
                        metadata={
                            "file_type": "pptx",
                            "slide_number": slide_index,
                            "filename": filename,
                        },
                        token_count=token_count,
                        chunk_index=len(chunks),
                    )
                )

        if not chunks:
            from core.config import settings

            if settings.LLAMA_CLOUD_API_KEY:
                if llamaparse_fallback_total:
                    llamaparse_fallback_total.labels("pptx").inc()
                return LlamaParseProcessor(file_type="pptx").process(content, filename)

        logger.info(f"[PPTXProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens")
        return ProcessedDocument(chunks=chunks, file_type="pptx", total_tokens=total_tokens)


# =============================================================================
# EMAIL PROCESSOR
# =============================================================================

class EmailProcessor(BaseProcessor):
    """Processor for .eml and .msg email files."""

    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        ext = os.path.splitext(filename)[1].lower()
        text = ""

        if ext == ".eml":
            try:
                from email import policy
                from email.parser import BytesParser
            except ImportError as exc:
                logger.error(f"[EmailProcessor] Missing dependency: {exc}")
                return ProcessedDocument(
                    chunks=[],
                    file_type="unsupported",
                    metadata={"unsupported_reason": "missing_dependency"},
                )

            msg = BytesParser(policy=policy.default).parsebytes(content)
            text = self._email_to_text(msg)
        elif ext == ".msg":
            try:
                import extract_msg
            except ImportError as exc:
                logger.error(f"[EmailProcessor] Missing dependency: {exc}")
                return ProcessedDocument(
                    chunks=[],
                    file_type="unsupported",
                    metadata={"unsupported_reason": "missing_dependency"},
                )

            with tempfile.NamedTemporaryFile(delete=False, suffix=".msg") as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            try:
                msg = extract_msg.Message(tmp_path)
                msg.process()
                text = self._msg_to_text(msg)
            finally:
                try:
                    os.remove(tmp_path)
                except Exception:
                    pass
        else:
            return ProcessedDocument(chunks=[], file_type="email")

        if not text.strip():
            from core.config import settings

            if settings.LLAMA_CLOUD_API_KEY:
                if llamaparse_fallback_total:
                    llamaparse_fallback_total.labels("email").inc()
                return LlamaParseProcessor(file_type="email").process(content, filename)
            return ProcessedDocument(chunks=[], file_type="email")

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=150,
            separators=["\n\n", "\n", ". ", " ", ""],
        )

        chunks = []
        total_tokens = 0
        for i, chunk_text in enumerate(splitter.split_text(text)):
            token_count = self.count_tokens(chunk_text)
            total_tokens += token_count
            chunks.append(
                ProcessedChunk(
                    content=chunk_text,
                    metadata={"file_type": "email", "filename": filename},
                    token_count=token_count,
                    chunk_index=i,
                )
            )

        logger.info(f"[EmailProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens")
        return ProcessedDocument(chunks=chunks, file_type="email", total_tokens=total_tokens)

    def _email_to_text(self, msg) -> str:
        parts = [
            f"Subject: {msg.get('subject', '')}",
            f"From: {msg.get('from', '')}",
            f"To: {msg.get('to', '')}",
            f"Date: {msg.get('date', '')}",
        ]

        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                if content_type == "text/plain":
                    body = part.get_content()
                    break
                if content_type == "text/html" and not body:
                    body = part.get_content()
        else:
            body = msg.get_content()

        if body:
            body_text = body
            if "<html" in body.lower():
                try:
                    from bs4 import BeautifulSoup
                    body_text = BeautifulSoup(body, "html.parser").get_text(separator="\n")
                except Exception:
                    body_text = body
            parts.append("Body:\n" + body_text.strip())

        return "\n".join(p for p in parts if p)

    def _msg_to_text(self, msg) -> str:
        parts = [
            f"Subject: {getattr(msg, 'subject', '')}",
            f"From: {getattr(msg, 'sender', '')}",
            f"To: {getattr(msg, 'to', '')}",
            f"Date: {getattr(msg, 'date', '')}",
        ]
        body = getattr(msg, "body", "") or ""
        if body:
            parts.append("Body:\n" + body.strip())
        return "\n".join(p for p in parts if p)


# =============================================================================
# LLAMAPARSE PROCESSOR
# =============================================================================

class LlamaParseProcessor(BaseProcessor):
    """Generic LlamaParse processor for legacy, binary, or image files."""

    def __init__(self, file_type: str = "llamaparse"):
        self.file_type = file_type

    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        from core.config import settings

        if not settings.LLAMA_CLOUD_API_KEY:
            logger.warning(f"[LlamaParseProcessor] LlamaParse not configured for {safe_file_ref(filename=filename)}")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "llamaparse_unavailable"},
            )

        try:
            import nest_asyncio
            nest_asyncio.apply()
            from llama_parse import LlamaParse
        except ImportError as exc:
            logger.error(f"[LlamaParseProcessor] Missing dependency: {exc}")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "llamaparse_missing_dependency"},
            )

        suffix = os.path.splitext(filename)[1] or ".bin"
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tf:
            tf.write(content)
            tf_path = tf.name

        try:
            with LLAMAPARSE_LOCK:
                parser = LlamaParse(
                    api_key=settings.LLAMA_CLOUD_API_KEY,
                    result_type="markdown",
                    verbose=False,
                )
                documents = parser.load_data(tf_path)
        finally:
            try:
                os.remove(tf_path)
            except Exception:
                pass

        if not documents:
            return ProcessedDocument(chunks=[], file_type=self.file_type)

        full_text = "\n\n".join([doc.text for doc in documents if getattr(doc, "text", None)])
        if not full_text.strip():
            return ProcessedDocument(chunks=[], file_type=self.file_type)

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""],
        )
        raw_chunks = splitter.split_text(full_text)

        chunks = []
        total_tokens = 0
        for i, chunk_text in enumerate(raw_chunks):
            contextualized = f"[File: {filename}]\n{chunk_text}"
            token_count = self.count_tokens(contextualized)
            total_tokens += token_count
            chunks.append(
                ProcessedChunk(
                    content=contextualized,
                    metadata={
                        "file_type": self.file_type,
                        "parser": "llama_parse",
                        "filename": filename,
                    },
                    token_count=token_count,
                    chunk_index=i,
                )
            )

        return ProcessedDocument(
            chunks=chunks,
            file_type=self.file_type,
            total_tokens=total_tokens,
            metadata={"parser": "llama_parse"},
        )


# =============================================================================
# LEGACY OFFICE PROCESSOR
# =============================================================================

class LegacyOfficeProcessor(BaseProcessor):
    """Processor for legacy Office formats routed to LlamaParse."""

    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        ext = os.path.splitext(filename)[1].lower().lstrip(".") or "legacy_office"
        return LlamaParseProcessor(file_type=ext).process(content, filename)


# =============================================================================
# IMAGE PROCESSOR (ENHANCED WITH LOCAL TESSERACT)
# =============================================================================

class ImageProcessor(BaseProcessor):
    """
    Processor for images with 4-Tier Resilient Architecture.

    Tier 1 (Fast Local): Tesseract OCR - Direct image-to-text
    Tier 2 (Cloud OCR): LlamaParse - Better quality for complex layouts
    Tier 3 (Vision LLM): GPT-4o/Gemini - Semantic understanding for diagrams
    Tier 4 (Metadata): EXIF/filename fallback

    Handles: JPG, JPEG, PNG, TIFF, BMP, GIF

    This ensures images can be processed even without cloud API access,
    with cascading fallback through multiple tiers for best quality.
    """

    MIN_TOKENS_THRESHOLD = 20  # Below this, try cloud fallback

    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        """
        Process image with 4-tier cascading fallback.

        Flow:
        1. Try Tesseract OCR locally
        2. If low content or failed, try LlamaParse (if configured)
        3. If still low content, try Vision LLM (if enabled)
        4. If all else fails, return metadata-only result
        """
        from core.config import settings

        # =====================================================================
        # TIER 1: Tesseract OCR (Fast Local)
        # =====================================================================
        logger.info(f"[ImageProcessor] Tier 1: Trying Tesseract OCR for {safe_file_ref(filename=filename)}")
        try:
            result = self._process_with_tesseract(content, filename)
            if result and result.chunks and result.total_tokens >= self.MIN_TOKENS_THRESHOLD:
                logger.info(f"[ImageProcessor] Tier 1 success: {result.total_tokens} tokens from {safe_file_ref(filename=filename)}")
                if image_ocr_total:
                    image_ocr_total.labels("tier1_tesseract_success").inc()
                return result
            logger.warning(
                f"[ImageProcessor] Tesseract returned low content "
                f"({result.total_tokens if result else 0} tokens) for {filename}"
            )
        except OCRNotAvailableException as e:
            logger.warning(f"[ImageProcessor] Tesseract not available: {e}")
        except Exception as e:
            logger.error(f"[ImageProcessor] Tesseract failed for {safe_file_ref(filename=filename)}: {e}")

        # =====================================================================
        # TIER 2: LlamaParse (Cloud OCR Fallback)
        # =====================================================================
        if settings.LLAMA_CLOUD_API_KEY:
            can_execute, reason = LLAMAPARSE_CIRCUIT.can_execute()
            if can_execute:
                logger.info(f"[ImageProcessor] Tier 2: Trying LlamaParse for {safe_file_ref(filename=filename)}")
                try:
                    result = LlamaParseProcessor(file_type="image").process(content, filename)
                    if result and result.chunks:
                        LLAMAPARSE_CIRCUIT.record_success()
                        if image_ocr_total:
                            image_ocr_total.labels("tier2_llamaparse_success").inc()
                        logger.info(f"[ImageProcessor] Tier 2 LlamaParse success: {result.total_tokens} tokens from {safe_file_ref(filename=filename)}")
                        return result
                except Exception as e:
                    error_str = str(e).lower()
                    if "402" in error_str or "quota" in error_str or "payment" in error_str or "credits" in error_str:
                        LLAMAPARSE_CIRCUIT.record_failure("402")
                        logger.warning(f"[ImageProcessor] LlamaParse quota exceeded: {e}")
                    else:
                        LLAMAPARSE_CIRCUIT.record_failure("error")
                        logger.warning(f"[ImageProcessor] LlamaParse failed: {e}")
            else:
                logger.info(f"[ImageProcessor] Skipping LlamaParse ({reason}) for {safe_file_ref(filename=filename)}")
                if llamaparse_fallback_total:
                    llamaparse_fallback_total.labels(f"image_{reason}").inc()

        # =====================================================================
        # TIER 3: Vision LLM (Semantic Understanding)
        # =====================================================================
        if settings.VISION_LLM_ENABLED:
            result = self._process_with_vision_llm(content, filename)
            if result and result.chunks:
                logger.info(f"[ImageProcessor] Tier 3 Vision LLM success: {result.total_tokens} tokens from {safe_file_ref(filename=filename)}")
                if image_ocr_total:
                    image_ocr_total.labels("tier3_vision_llm_success").inc()
                return result

        # =====================================================================
        # TIER 4: Metadata Fallback
        # =====================================================================
        # Return empty if nothing worked
        if image_ocr_total:
            image_ocr_total.labels("all_tiers_failed").inc()
        logger.warning(f"[ImageProcessor] All tiers failed for {safe_file_ref(filename=filename)}")
        return ProcessedDocument(chunks=[], file_type="image")

    def _process_with_vision_llm(self, content: bytes, filename: str) -> ProcessedDocument:
        """
        Process image using Vision LLM for semantic understanding.

        Uses GPT-4o or Gemini Pro Vision to analyze diagrams, charts,
        and technical images, extracting semantic meaning rather than
        just text.

        Returns:
            ProcessedDocument with Vision LLM analysis
        """
        import asyncio
        import time
        from core.config import settings
        from services.secure_cleanup import SmartBuffer

        try:
            from core.metrics import vision_llm_total, vision_llm_duration, vision_llm_diagram_types
        except ImportError:
            vision_llm_total = None
            vision_llm_duration = None
            vision_llm_diagram_types = None

        # Check file size
        if len(content) > settings.VISION_LLM_MAX_IMAGE_SIZE:
            logger.warning(f"[ImageProcessor] Image too large for Vision LLM: {len(content)} bytes")
            if vision_llm_total:
                vision_llm_total.labels(provider=settings.VISION_LLM_PROVIDER, result="skipped_size").inc()
            return ProcessedDocument(chunks=[], file_type="image")

        # Use SmartBuffer for Ghost Protocol compliance
        with SmartBuffer(content, filename=filename) as buffer:
            try:
                # Get appropriate processor
                from services.vision.circuit import get_circuit_for_provider
                circuit = get_circuit_for_provider(settings.VISION_LLM_PROVIDER)

                can_execute, reason = circuit.can_execute()
                if not can_execute:
                    logger.info(f"[ImageProcessor] Skipping Vision LLM ({reason})")
                    if vision_llm_total:
                        vision_llm_total.labels(provider=settings.VISION_LLM_PROVIDER, result=f"skipped_{reason}").inc()
                    return ProcessedDocument(chunks=[], file_type="image")

                # Get processor based on provider
                if settings.VISION_LLM_PROVIDER == "gemini":
                    from services.vision.gemini_vision import GeminiVisionProcessor
                    processor = GeminiVisionProcessor()
                else:
                    from services.vision.openai_vision import GPT4VisionProcessor
                    processor = GPT4VisionProcessor()

                if not processor.is_available():
                    logger.warning(f"[ImageProcessor] Vision LLM not available: {settings.VISION_LLM_PROVIDER}")
                    if vision_llm_total:
                        vision_llm_total.labels(provider=settings.VISION_LLM_PROVIDER, result="not_available").inc()
                    return ProcessedDocument(chunks=[], file_type="image")

                # Run analysis
                start_time = time.time()

                # Run async in sync context
                loop = asyncio.new_event_loop()
                try:
                    vision_result = loop.run_until_complete(
                        processor.analyze_image(buffer.get_bytes(), filename)
                    )
                finally:
                    loop.close()

                duration = time.time() - start_time

                # Record metrics
                if vision_llm_total:
                    vision_llm_total.labels(provider=settings.VISION_LLM_PROVIDER, result="success").inc()
                if vision_llm_duration:
                    vision_llm_duration.labels(provider=settings.VISION_LLM_PROVIDER).observe(duration)
                if vision_llm_diagram_types and vision_result.diagram_type:
                    vision_llm_diagram_types.labels(diagram_type=vision_result.diagram_type.value).inc()

                circuit.record_success()

                # Convert to searchable text
                text = vision_result.to_searchable_text()

                if not text.strip():
                    return ProcessedDocument(chunks=[], file_type="image")

                # Chunk the vision analysis text
                splitter = RecursiveCharacterTextSplitter(
                    chunk_size=1000,
                    chunk_overlap=200,
                    separators=["\n\n", "\n", ". ", " ", ""]
                )
                raw_chunks = splitter.split_text(text)

                chunks = []
                total_tokens = 0
                for i, chunk_text in enumerate(raw_chunks):
                    contextualized = f"[File: {filename}]\n{chunk_text}"
                    token_count = self.count_tokens(contextualized)
                    total_tokens += token_count

                    chunks.append(ProcessedChunk(
                        content=contextualized,
                        metadata={
                            "file_type": "image",
                            "parser": "vision_llm",
                            "vision_provider": settings.VISION_LLM_PROVIDER,
                            "diagram_type": vision_result.diagram_type.value if vision_result.diagram_type else None,
                            "confidence": vision_result.confidence,
                            "filename": filename,
                        },
                        token_count=token_count,
                        chunk_index=i
                    ))

                logger.info(
                    f"[ImageProcessor] Vision LLM: {safe_file_ref(filename=filename)}: "
                    f"{len(chunks)} chunks, {total_tokens} tokens, "
                    f"type={vision_result.diagram_type.value if vision_result.diagram_type else 'unknown'}"
                )

                return ProcessedDocument(
                    chunks=chunks,
                    file_type="image",
                    total_tokens=total_tokens,
                    metadata={
                        "parser": "vision_llm",
                        "vision_provider": settings.VISION_LLM_PROVIDER,
                        "diagram_type": vision_result.diagram_type.value if vision_result.diagram_type else None,
                    }
                )

            except Exception as e:
                error_str = str(e).lower()
                if "402" in error_str or "quota" in error_str or "rate" in error_str:
                    circuit.record_failure("quota")
                else:
                    circuit.record_failure("error")

                logger.error(f"[ImageProcessor] Vision LLM failed for {safe_file_ref(filename=filename)}: {e}")
                if vision_llm_total:
                    vision_llm_total.labels(provider=settings.VISION_LLM_PROVIDER, result="failure").inc()
                return ProcessedDocument(chunks=[], file_type="image")
    
    def _process_with_tesseract(self, content: bytes, filename: str) -> ProcessedDocument:
        """
        Run Tesseract OCR directly on image.
        
        Includes preprocessing for better OCR accuracy:
        - Color mode conversion
        - Image enhancement for low-quality images
        
        Returns:
            ProcessedDocument with OCR'd text
        """
        try:
            import pytesseract
            from PIL import Image, ImageEnhance, ImageFilter
        except ImportError as e:
            raise OCRNotAvailableException(f"OCR dependencies not installed: {e}")
        
        # Load image
        try:
            image = Image.open(io.BytesIO(content))
            
            # Get original dimensions for logging
            original_size = f"{image.width}x{image.height}"
            
            # Convert to RGB if necessary (handles RGBA, P, LA, L modes)
            if image.mode in ('RGBA', 'P', 'LA'):
                # Create white background for transparent images
                background = Image.new('RGB', image.size, (255, 255, 255))
                if image.mode == 'P':
                    image = image.convert('RGBA')
                if image.mode in ('RGBA', 'LA'):
                    background.paste(image, mask=image.split()[-1])
                else:
                    background.paste(image)
                image = background
            elif image.mode == 'L':
                image = image.convert('RGB')
            elif image.mode not in ('RGB',):
                image = image.convert('RGB')
            
            # Apply light preprocessing for better OCR
            # Sharpen slightly to improve text clarity
            image = image.filter(ImageFilter.SHARPEN)
            
            # Enhance contrast slightly
            enhancer = ImageEnhance.Contrast(image)
            image = enhancer.enhance(1.2)
            
            # Run OCR with language detection
            # Try multiple languages for better accuracy
            try:
                text = pytesseract.image_to_string(image, lang='eng+tur')
            except Exception:
                # Fallback to English only if multi-lang fails
                text = pytesseract.image_to_string(image, lang='eng')
            
            logger.debug(f"[ImageProcessor] Processed {safe_file_ref(filename=filename)} ({original_size}), extracted {len(text)} chars")
            
        except Exception as e:
            logger.error(f"[ImageProcessor] Failed to process image {safe_file_ref(filename=filename)}: {e}")
            raise
        
        if not text.strip():
            return ProcessedDocument(chunks=[], file_type="image", total_tokens=0)
        
        # Clean up OCR text
        text = self._clean_ocr_text(text)
        
        if not text.strip():
            return ProcessedDocument(chunks=[], file_type="image", total_tokens=0)
        
        # Chunk the OCR text
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        raw_chunks = splitter.split_text(text)
        
        chunks = []
        total_tokens = 0
        for i, chunk_text in enumerate(raw_chunks):
            contextualized = f"[File: {filename}]\n{chunk_text}"
            token_count = self.count_tokens(contextualized)
            total_tokens += token_count
            
            chunks.append(ProcessedChunk(
                content=contextualized,
                metadata={
                    "file_type": "image",
                    "parser": "tesseract_ocr",
                    "filename": filename,
                },
                token_count=token_count,
                chunk_index=i
            ))
        
        logger.info(f"[ImageProcessor] Tesseract: {safe_file_ref(filename=filename)}: {len(chunks)} chunks, {total_tokens} tokens")
        return ProcessedDocument(
            chunks=chunks,
            file_type="image",
            total_tokens=total_tokens,
            metadata={"parser": "tesseract_ocr"}
        )
    
    def _clean_ocr_text(self, text: str) -> str:
        """Clean up common OCR artifacts and noise."""
        # Remove excessive whitespace
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r' {2,}', ' ', text)
        
        # Remove common OCR noise patterns
        lines = text.split('\n')
        cleaned_lines = []
        for line in lines:
            line = line.strip()
            # Skip empty lines at the start
            if not cleaned_lines and not line:
                continue
            # Skip lines that are just noise (single characters, numbers only, etc.)
            if len(line) <= 2 and not line.isalnum():
                continue
            # Skip lines that are mostly special characters
            if line and sum(1 for c in line if c.isalnum()) / len(line) < 0.3:
                continue
            cleaned_lines.append(line)
        
        return '\n'.join(cleaned_lines)


# =============================================================================
# PLAIN TEXT PROCESSOR
# =============================================================================

class PlainTextProcessor(BaseProcessor):
    """Processor for plain text files."""
    
    def process(self, content: bytes, filename: str) -> ProcessedDocument:
        """Process plain text with simple chunking."""
        try:
            text = content.decode("utf-8")
        except UnicodeDecodeError:
            text = content.decode("utf-8", errors="replace")
        
        # Remove null bytes that cannot be stored in Postgres text
        text = text.replace("\x00", "")
        
        if not text.strip():
            return ProcessedDocument(chunks=[], file_type="text")
        
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=150,
            separators=["\n\n", "\n", ". ", " ", ""]
        )
        raw_chunks = splitter.split_text(text)
        
        chunks = []
        total_tokens = 0
        for i, chunk_text in enumerate(raw_chunks):
            token_count = self.count_tokens(chunk_text)
            total_tokens += token_count
            chunks.append(ProcessedChunk(
                content=chunk_text,
                metadata={"file_type": "text", "filename": filename},
                token_count=token_count,
                chunk_index=i
            ))
        
        logger.info(f"[PlainTextProcessor] {safe_file_ref(filename=filename)}: {len(chunks)} chunks")
        return ProcessedDocument(chunks=chunks, file_type="text", total_tokens=total_tokens)


# =============================================================================
# DOCUMENT PROCESSOR FACTORY
# =============================================================================

class DocumentProcessorFactory:
    """
    Factory class that selects the appropriate processor based on file type.
    
    Usage:
        result = DocumentProcessorFactory.process(file_path, filename)
        chunks = result.chunks  # List[ProcessedChunk]
    """
    
    # Extension to processor mapping
    PROCESSOR_MAP = {
        # Code files
        ".py": CodeProcessor,
        ".js": CodeProcessor,
        ".jsx": CodeProcessor,
        ".ts": CodeProcessor,
        ".tsx": CodeProcessor,
        ".java": CodeProcessor,
        ".go": CodeProcessor,
        ".cpp": CodeProcessor,
        ".c": CodeProcessor,
        ".cs": CodeProcessor,
        ".rb": CodeProcessor,
        ".php": CodeProcessor,
        ".rs": CodeProcessor,
        ".scala": CodeProcessor,
        ".swift": CodeProcessor,
        ".kt": CodeProcessor,
        ".json": CodeProcessor,
        ".yaml": CodeProcessor,
        ".yml": CodeProcessor,
        ".toml": CodeProcessor,
        ".xml": CodeProcessor,
        ".ini": CodeProcessor,
        ".conf": CodeProcessor,
        ".config": CodeProcessor,
        ".sh": CodeProcessor,
        ".dockerfile": CodeProcessor,
        ".html": HTMLProcessor,
        ".css": CodeProcessor,
        ".sql": CodeProcessor,
        
        # Markdown
        ".md": MarkdownProcessor,
        ".markdown": MarkdownProcessor,
        
        # PDF
        ".pdf": PDFProcessor,
        
        # Word documents
        ".docx": DocxProcessor,
        ".doc": LegacyOfficeProcessor,

        # Spreadsheets
        ".csv": CSVProcessor,
        ".tsv": CSVProcessor,
        ".xlsx": ExcelProcessor,
        ".xls": LegacyOfficeProcessor,

        # Presentations
        ".pptx": PPTXProcessor,
        ".ppt": LegacyOfficeProcessor,

        # Rich text / email
        ".rtf": LegacyOfficeProcessor,
        ".msg": EmailProcessor,
        ".eml": EmailProcessor,

        # Images (OCR)
        ".jpg": ImageProcessor,
        ".jpeg": ImageProcessor,
        ".png": ImageProcessor,
        ".tiff": ImageProcessor,
        ".bmp": ImageProcessor,
        
        # Plain text
        ".txt": PlainTextProcessor,
        ".log": PlainTextProcessor,
    }
    
    # MIME type fallback mapping
    MIME_MAP = {
        "application/pdf": PDFProcessor,
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": DocxProcessor,
        "application/msword": LegacyOfficeProcessor,
        "text/markdown": MarkdownProcessor,
        "text/plain": PlainTextProcessor,
        "text/html": HTMLProcessor,
        "text/csv": CSVProcessor,
        "application/json": CodeProcessor,
        "application/xml": CodeProcessor,
        "text/xml": CodeProcessor,
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": ExcelProcessor,
        "application/vnd.ms-excel": LegacyOfficeProcessor,
        "application/vnd.openxmlformats-officedocument.presentationml.presentation": PPTXProcessor,
        "application/vnd.ms-powerpoint": LegacyOfficeProcessor,
        "message/rfc822": EmailProcessor,
        "application/vnd.ms-outlook": EmailProcessor,
        "application/rtf": LegacyOfficeProcessor,
        "image/jpeg": ImageProcessor,
        "image/jpg": ImageProcessor,
        "image/png": ImageProcessor,
        "image/tiff": ImageProcessor,
        "image/bmp": ImageProcessor,
    }

    # Explicitly unsupported (binary) extensions to avoid unsafe parsing
    UNSUPPORTED_EXTENSIONS = {
        ".numbers",
        ".key",
    }

    TEXT_MIME_TYPES = {
        "text/plain",
        "text/markdown",
        "text/csv",
        "text/html",
        "text/xml",
        "application/json",
        "application/xml",
        "application/xhtml+xml",
        "application/x-yaml",
    }

    @staticmethod
    def _looks_like_binary(content: bytes) -> bool:
        if not content:
            return False
        if b"\x00" in content:
            return True
        sample = content[:2048]
        if not sample:
            return False
        non_text = 0
        for byte in sample:
            if byte < 9 or (14 <= byte < 32) or byte == 127:
                non_text += 1
        return (non_text / len(sample)) > 0.30
    
    @classmethod
    def process(
        cls,
        file_path: str = None,
        filename: str = None,
        content: bytes = None,
        mime_type: str = None
    ) -> ProcessedDocument:
        """
        Process a document using the appropriate processor.
        
        Args:
            file_path: Path to file on disk (optional if content provided)
            filename: Original filename for extension detection
            content: Raw file bytes (optional if file_path provided)
            mime_type: MIME type for fallback detection
        
        Returns:
            ProcessedDocument with chunks and metadata
        """
        # Load content if file_path provided
        if file_path and not content:
            with open(file_path, "rb") as f:
                content = f.read()
        
        if not content:
            logger.error("[Factory] No content to process")
            return ProcessedDocument(chunks=[], file_type="unknown")
        
        # Determine processor
        filename = filename or (os.path.basename(file_path) if file_path else "unknown")
        ext = os.path.splitext(filename)[1].lower()
        if not ext and filename and filename.lower() == "dockerfile":
            ext = ".dockerfile"
        
        processor_class = cls.PROCESSOR_MAP.get(ext)
        
        if not processor_class and mime_type:
            processor_class = cls.MIME_MAP.get(mime_type)
        
        if not processor_class:
            if ext in cls.UNSUPPORTED_EXTENSIONS:
                logger.warning(f"[Factory] Unsupported file type {ext}, skipping parse")
                return ProcessedDocument(
                    chunks=[],
                    file_type="unsupported",
                    metadata={"unsupported_reason": "unsupported_extension"}
                )
            if mime_type and (mime_type.startswith("text/") or mime_type in cls.TEXT_MIME_TYPES):
                processor_class = PlainTextProcessor
            elif not cls._looks_like_binary(content):
                processor_class = PlainTextProcessor
            else:
                logger.warning(f"[Factory] Binary content detected for {ext}, skipping parse")
                return ProcessedDocument(
                    chunks=[],
                    file_type="unsupported",
                    metadata={"unsupported_reason": "binary_content"}
                )

        if processor_class is PlainTextProcessor and cls._looks_like_binary(content):
            logger.warning(f"[Factory] Binary content detected for {ext}, skipping plain text parse")
            return ProcessedDocument(
                chunks=[],
                file_type="unsupported",
                metadata={"unsupported_reason": "binary_content"}
            )
        
        # Process
        processor = processor_class()
        result = processor.process(content, filename)
        
        logger.info(f"[Factory] Processed {safe_file_ref(filename=filename)}: {len(result.chunks)} chunks, type={result.file_type}")
        return result
    
    @classmethod
    def process_web_content(cls, html_content: str, url: str) -> ProcessedDocument:
        """
        Special method for processing web content (already converted to text/markdown).
        
        Args:
            html_content: Text content extracted from web page
            url: Source URL for metadata
        
        Returns:
            ProcessedDocument with chunks
        """
        processor = MarkdownProcessor()
        content_bytes = html_content.encode("utf-8")
        result = processor.process(content_bytes, url)
        
        # Add source URL to all chunk metadata
        for chunk in result.chunks:
            chunk.metadata["source_url"] = url
        
        return result


# =============================================================================
# ROUTING HELPERS
# =============================================================================

def route_file(file_size: int, extension: str) -> str:
    """Route structured files based on size limits (no heavy queue yet)."""
    from core.config import settings

    if extension in {".csv", ".tsv", ".xlsx"}:
        if file_size > settings.MAX_STRUCTURED_FILE_SIZE:
            return "skipped_file_too_large"
    return "queues.parsing"


# =============================================================================
# LEGACY COMPATIBILITY - DocumentParser
# =============================================================================

class DocumentParser:
    """
    Legacy compatibility class.
    
    Maps to the old extract_text interface for backwards compatibility.
    New code should use DocumentProcessorFactory directly.
    """
    
    SUPPORTED_FORMATS = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/msword': 'doc',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'application/vnd.ms-excel': 'xls',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'application/vnd.ms-powerpoint': 'ppt',
        'text/plain': 'text',
        'text/markdown': 'text',
        'text/csv': 'text',
        'text/html': 'text',
        'application/xml': 'text',
        'text/xml': 'text',
        'message/rfc822': 'email',
        'application/vnd.ms-outlook': 'email',
        'application/rtf': 'rtf',
    }
    
    @staticmethod
    def extract_text(file_content: bytes, mime_type: str) -> str:
        """Extract text from file bytes (legacy method)."""
        result = DocumentProcessorFactory.process(
            content=file_content,
            filename="document",
            mime_type=mime_type
        )
        # Combine all chunks into single text
        return "\n\n".join(chunk.content for chunk in result.chunks)
    
    @staticmethod
    def parse_file(file_path: str, filename: str = None) -> str:
        """Parse file and return combined text (legacy method)."""
        result = DocumentProcessorFactory.process(
            file_path=file_path,
            filename=filename or os.path.basename(file_path)
        )
        return "\n\n".join(chunk.content for chunk in result.chunks)
    
    @staticmethod
    def is_supported(mime_type: str) -> bool:
        """Check if MIME type is supported."""
        return mime_type.lower().strip() in DocumentParser.SUPPORTED_FORMATS
